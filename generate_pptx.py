"""Convert rail-pitch-deck.html to rail-pitch-deck.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import re, base64, io, os

# ── colours ──────────────────────────────────────────────────────────────────
RED   = RGBColor(0xFF, 0x2E, 0x01)
GREEN = RGBColor(0x00, 0xC8, 0x53)
BLACK = RGBColor(0x08, 0x08, 0x08)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF6, 0xF5, 0xF2)
GRAY4 = RGBColor(0x96, 0x93, 0x8C)
GRAY6 = RGBColor(0x56, 0x53, 0x4E)

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── helpers ──────────────────────────────────────────────────────────────────
def add_slide(bg: RGBColor):
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide

def box(slide, x, y, w, h, text, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Plus Jakarta Sans"
    return txBox

def rect(slide, x, y, w, h, fill: RGBColor, radius=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def eyebrow(slide, text, y, dark=False):
    box(slide, Inches(1), y, Inches(8), Inches(0.35),
        text.upper(), size=9, bold=True,
        color=RED)

def rule_line(slide, y):
    rect(slide, Inches(1), y, Inches(0.5), Pt(2), RED)

def bullet_point(slide, text, y, color=WHITE):
    # red dot
    dot = slide.shapes.add_shape(1, Inches(1), y + Inches(0.08), Pt(6), Pt(6))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
    box(slide, Inches(1.2), y, Inches(9), Inches(0.4), text, size=13, color=color)

def callout_box(slide, text, y, bg_rgb, border_rgb, text_color=WHITE, height=Inches(0.7)):
    r = rect(slide, Inches(1), y, Inches(8), height, bg_rgb)
    # left accent
    rect(slide, Inches(1), y, Pt(3), height, border_rgb)
    box(slide, Inches(1.25), y + Inches(0.1), Inches(7.5), height - Inches(0.2),
        text, size=13, bold=True, color=text_color, wrap=True)

def stat_card(slide, x, y, w, h, num, label, bg=WHITE, num_color=BLACK, label_color=GRAY6):
    rect(slide, x, y, w, h, bg)
    box(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.55),
        num, size=26, bold=True, color=num_color)
    box(slide, x + Inches(0.15), y + Inches(0.65), w - Inches(0.3), h - Inches(0.75),
        label, size=10, color=label_color, wrap=True)

def stream_row(slide, y, icon, title, desc):
    box(slide, Inches(1), y, Inches(0.5), Inches(0.5), icon, size=18)
    box(slide, Inches(1.6), y, Inches(8), Inches(0.28), title, size=13, bold=True, color=WHITE)
    box(slide, Inches(1.6), y + Inches(0.28), Inches(8), Inches(0.28), desc, size=11, color=GRAY4)

# ── extract base64 image from HTML ───────────────────────────────────────────
html_path = os.path.join(os.path.dirname(__file__), "rail-pitch-deck.html")
with open(html_path, "rb") as f:
    raw = f.read().decode("utf-8", errors="replace")

b64_match = re.search(r'data:image/([^;]+);base64,([A-Za-z0-9+/=]+)', raw)
hero_img_stream = None
if b64_match:
    img_data = base64.b64decode(b64_match.group(2))
    hero_img_stream = io.BytesIO(img_data)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — INTRO
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(BLACK)
box(s, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
    "Rail", size=22, bold=True, color=WHITE)
box(s, Inches(1), Inches(1.3), Inches(7), Inches(1.8),
    "Your money account that\ninvests for you — automatically.",
    size=36, bold=True, color=WHITE)
# "invests for you" accent — overlay red text on same area
box(s, Inches(1), Inches(1.3), Inches(7), Inches(1.8),
    "\n\ninvests for you", size=36, bold=True, color=RED)
box(s, Inches(1), Inches(3.2), Inches(7), Inches(1.2),
    "Spend normally. Rail allocates a portion of every deposit toward long-term\n"
    "investments in the background. No dashboards. No decisions. No discipline required.",
    size=13, color=RGBColor(0xAA, 0xAA, 0xAA))
box(s, Inches(1), Inches(4.6), Inches(5), Inches(0.4),
    "Solana Mobile · Monolith Hackathon 2026", size=10,
    color=RGBColor(0x88, 0x88, 0x88))
if hero_img_stream:
    try:
        s.shapes.add_picture(hero_img_stream, Inches(8.5), Inches(0.3),
                             height=Inches(6.8))
    except Exception:
        pass
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "01 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(BLACK)
eyebrow(s, "The Problem", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(1.4),
    "Millions work for years but never build\nmeaningful investments.",
    size=28, bold=True, color=WHITE)
pts = [
    "Saving requires constant discipline most people don't sustain",
    "Investing requires decisions people avoid making",
    "Most financial apps rely on users to stay consistent — they don't",
]
for i, pt in enumerate(pts):
    bullet_point(s, pt, Inches(2.7 + i * 0.55))
callout_box(s, "If investing is not automatic, most people never do it consistently.",
            Inches(4.6), RGBColor(0x1A, 0x08, 0x05), RED, height=Inches(0.65))
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "02 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VALUE PROP
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(CREAM)
eyebrow(s, "Value Proposition", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(1.4),
    "A primary money account where spending feels normal\nand investing happens automatically.",
    size=26, bold=True, color=BLACK)
cards = [
    ("💸", "No market tracking", "Users never need to watch charts or manage dashboards."),
    ("⚡", "Auto-allocation", "Every deposit is automatically split between spending and investing."),
    ("🔄", "Zero behavior change", "Users spend normally. Rail builds wealth in the background."),
]
card_w = Inches(3.6)
for i, (icon, title, body) in enumerate(cards):
    cx = Inches(1 + i * 3.8)
    cy = Inches(2.8)
    rect(s, cx, cy, card_w, Inches(2.2), WHITE)
    box(s, cx + Inches(0.2), cy + Inches(0.15), card_w - Inches(0.4), Inches(0.5),
        icon, size=18)
    box(s, cx + Inches(0.2), cy + Inches(0.7), card_w - Inches(0.4), Inches(0.4),
        title, size=13, bold=True, color=BLACK)
    box(s, cx + Inches(0.2), cy + Inches(1.1), card_w - Inches(0.4), Inches(0.9),
        body, size=11, color=GRAY6, wrap=True)
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "03 / 11", size=9, color=GRAY4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(BLACK)
eyebrow(s, "How It Works", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(0.7),
    "Three steps. That's it.", size=28, bold=True, color=WHITE)
steps = [
    ("01", "Deposit funds into Rail",
     "Via bank transfer, crypto, or card — money lands in your Rail account."),
    ("02", "Rail splits automatically",
     "A set percentage is allocated to your investment stash — no action needed."),
    ("03", "Investments execute in the background",
     "While you spend normally through your Rail card, your portfolio grows."),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.1 + i * 1.1)
    box(s, Inches(1), y, Inches(0.7), Inches(0.5), num, size=22, bold=True, color=RED)
    box(s, Inches(1.8), y, Inches(8), Inches(0.3), title, size=14, bold=True, color=WHITE)
    box(s, Inches(1.8), y + Inches(0.3), Inches(8), Inches(0.35), desc, size=12, color=GRAY4)
callout_box(s,
    "Users build long-term financial progress without changing their spending behavior.",
    Inches(5.5), RGBColor(0x0A, 0x0A, 0x0A), GREEN,
    text_color=RGBColor(0xCC, 0xFF, 0xCC), height=Inches(0.65))
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "04 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DEMO
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x0C, 0x0C, 0x0C))
eyebrow(s, "Product Demo", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(0.6),
    "See it in action.", size=28, bold=True, color=WHITE)
# flow steps
flow = [("💳","Deposit funds"), ("⚡","Rail allocates"), ("📈","Investment grows"), ("🛍️","Spend normally")]
for i, (icon, label) in enumerate(flow):
    fx = Inches(1 + i * 2.8)
    rect(s, fx, Inches(2.0), Inches(2.4), Inches(0.6), RGBColor(0x1A, 0x1A, 0x1A))
    box(s, fx + Inches(0.1), Inches(2.05), Inches(2.2), Inches(0.5),
        f"{icon}  {label}", size=12, bold=True, color=WHITE)
    if i < 3:
        box(s, fx + Inches(2.45), Inches(2.15), Inches(0.3), Inches(0.3),
            "→", size=14, color=GRAY4)
# metrics
metrics = [
    ("Spending Stash", "$840.00", "Available to spend", WHITE),
    ("Investment Stash", "$312.50", "+4.2% this month", GREEN),
    ("Auto-invested", "$1,152.50", "Total, all time", GREEN),
]
for i, (lbl, val, sub, vc) in enumerate(metrics):
    mx = Inches(1 + i * 3.9)
    rect(s, mx, Inches(2.9), Inches(3.6), Inches(1.6), RGBColor(0x14, 0x14, 0x14))
    box(s, mx + Inches(0.2), Inches(3.0), Inches(3.2), Inches(0.3),
        lbl.upper(), size=9, bold=True, color=GRAY4)
    box(s, mx + Inches(0.2), Inches(3.3), Inches(3.2), Inches(0.5),
        val, size=22, bold=True, color=vc)
    box(s, mx + Inches(0.2), Inches(3.8), Inches(3.2), Inches(0.3),
        sub, size=10, color=GRAY4)
callout_box(s,
    "💡  Every time money enters Rail, investing happens automatically. No manual action required.",
    Inches(4.8), RGBColor(0x1A, 0x08, 0x05), RED, height=Inches(0.65))
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "05 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(CREAM)
eyebrow(s, "Market", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(1.0),
    "People whose income already flows digitally.",
    size=26, bold=True, color=BLACK)
tags = ["Remote workers", "Freelancers", "Creators", "Global professionals", "Crypto-native earners"]
tx = Inches(1)
for tag in tags:
    tw = Inches(len(tag) * 0.12 + 0.6)
    rect(s, tx, Inches(2.3), tw, Inches(0.38), WHITE)
    box(s, tx + Inches(0.15), Inches(2.33), tw - Inches(0.2), Inches(0.32),
        tag, size=11, bold=True, color=BLACK)
    tx += tw + Inches(0.15)
stats = [
    ("$200B+", "Stablecoin transaction volume in 2024 — growing fast"),
    ("35M+",   "Digital nomads and remote workers globally"),
    ("↑3×",    "Stablecoin adoption growth over the last 2 years"),
    ("$0",     "Automated investing tools built for this audience"),
]
for i, (num, lbl) in enumerate(stats):
    col = i % 2; row = i // 2
    sx = Inches(1 + col * 4.2)
    sy = Inches(3.0 + row * 1.5)
    stat_card(s, sx, sy, Inches(3.9), Inches(1.3), num, lbl)
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "06 / 11", size=9, color=GRAY4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(BLACK)
eyebrow(s, "Business Model", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(0.6),
    "Three clear revenue streams.", size=28, bold=True, color=WHITE)
streams = [
    ("💳", "Card Interchange", "Revenue every time a user spends through their Rail card."),
    ("📊", "Asset Management Fees", "Small percentage on assets under management — scales with user growth."),
    ("⭐", "Rail+ Premium", "Advanced financial tools, higher allocation limits, priority support."),
]
for i, (icon, title, desc) in enumerate(streams):
    ry = Inches(2.1 + i * 1.0)
    rect(s, Inches(1), ry, Inches(9), Inches(0.85), RGBColor(0x14, 0x14, 0x14))
    stream_row(s, ry + Inches(0.15), icon, title, desc)
callout_box(s,
    "10,000 users × $1,000 invested = $10M AUM. A 0.5% management fee = $50K ARR — before interchange or premium revenue.",
    Inches(5.4), RGBColor(0x03, 0x10, 0x07), GREEN,
    text_color=RGBColor(0xCC, 0xFF, 0xCC), height=Inches(0.75))
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "07 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TRACTION
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(CREAM)
eyebrow(s, "Traction", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(0.6),
    "Early but real.", size=28, bold=True, color=BLACK)
traction = [
    ("30+", "Early testers ready to use the product on launch"),
    ("✓",   "Full product flow built and running on TestFlight"),
    ("✓",   "Investment automation system live end-to-end"),
]
for i, (num, lbl) in enumerate(traction):
    cx = Inches(1 + i * 3.9)
    rect(s, cx, Inches(2.1), Inches(3.6), Inches(2.0), WHITE)
    box(s, cx + Inches(0.2), Inches(2.2), Inches(3.2), Inches(0.7),
        num, size=32, bold=True, color=RED if num == "30+" else BLACK)
    box(s, cx + Inches(0.2), Inches(2.95), Inches(3.2), Inches(0.9),
        lbl, size=11, color=GRAY6, wrap=True)
box(s, Inches(1), Inches(4.4), Inches(10), Inches(0.7),
    "Infrastructure partnerships in progress. Initial focus: controlled beta to validate user behavior and retention before scaling acquisition.",
    size=11, color=GRAY4, wrap=True)
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "08 / 11", size=9, color=GRAY4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(BLACK)
eyebrow(s, "Roadmap", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(0.6),
    "Where we're going.", size=28, bold=True, color=WHITE)
milestones = [
    ("Q1 2026 — Now", "Private beta launch with early testers. Core auto-invest flow live.", True),
    ("Q2 2026",       "Public product launch. Rail card rollout. User acquisition begins.", False),
    ("Q3 2026",       "Expanded automated investing strategies. Retention optimization.", False),
    ("Q4 2026",       "Lending, smart allocations, and advanced financial automation.", False),
]
for i, (quarter, milestone, is_now) in enumerate(milestones):
    my = Inches(2.1 + i * 1.1)
    dot_color = RED if is_now else RGBColor(0x44, 0x44, 0x44)
    dot = s.shapes.add_shape(1, Inches(1), my + Inches(0.1), Pt(10), Pt(10))
    dot.fill.solid(); dot.fill.fore_color.rgb = dot_color; dot.line.fill.background()
    box(s, Inches(1.4), my, Inches(8), Inches(0.3),
        quarter, size=10, bold=True, color=RED if is_now else GRAY4)
    box(s, Inches(1.4), my + Inches(0.3), Inches(8), Inches(0.4),
        milestone, size=13, bold=is_now, color=WHITE if is_now else RGBColor(0xAA, 0xAA, 0xAA))
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "09 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(CREAM)
eyebrow(s, "Team", Inches(0.6))
box(s, Inches(1), Inches(1.1), Inches(9), Inches(0.8),
    "Built by someone who lived this problem.", size=26, bold=True, color=BLACK)
# card
rect(s, Inches(1), Inches(2.1), Inches(7), Inches(4.5), WHITE)
box(s, Inches(1.3), Inches(2.3), Inches(5), Inches(0.5),
    "Tobi", size=20, bold=True, color=BLACK)
box(s, Inches(1.3), Inches(2.8), Inches(5), Inches(0.3),
    "FOUNDER & BUILDER", size=10, bold=True, color=GRAY4)
facts = [
    "Several years building products and shipping to real users",
    "Multiple hackathon participations — knows how to ship fast",
    "Rail was inspired by growing up in a household where financial stability was difficult despite consistent work",
]
for i, fact in enumerate(facts):
    bullet_point(s, fact, Inches(3.3 + i * 0.55), color=GRAY6)
callout_box(s,
    '"Build a system where financial progress happens automatically — for everyone, not just the disciplined."',
    Inches(5.3), RGBColor(0xF0, 0xEC, 0xE8), RED, text_color=BLACK, height=Inches(0.75))
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "10 / 11", size=9, color=GRAY4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CTA
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(BLACK)
box(s, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
    "Rail", size=22, bold=True, color=WHITE)
box(s, Inches(1.5), Inches(1.5), Inches(10), Inches(2.0),
    "Financial progress\nshould be automatic.",
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(1.5), Inches(3.7), Inches(10), Inches(0.6),
    "Join the waitlist. Test the product. Tell us what you think.",
    size=14, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)
# buttons (visual only)
rect(s, Inches(4.2), Inches(4.5), Inches(2.2), Inches(0.55), RED)
box(s, Inches(4.2), Inches(4.52), Inches(2.2), Inches(0.5),
    "Join the waitlist →", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, Inches(6.7), Inches(4.5), Inches(1.8), Inches(0.55), RGBColor(0x22, 0x22, 0x22))
box(s, Inches(6.7), Inches(4.52), Inches(1.8), Inches(0.5),
    "Try the app", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(5.5), Inches(5.3), Inches(2.5), Inches(0.35),
    "rail.app", size=11, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)
box(s, Inches(12), Inches(7), Inches(1), Inches(0.4),
    "11 / 11", size=9, color=RGBColor(0x44, 0x44, 0x44))

# ── save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "rail-pitch-deck.pptx")
prs.save(out)
print(f"Saved → {out}")
